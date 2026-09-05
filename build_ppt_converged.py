"""Converged 5-page deck: the explainer's plain-language substance, condensed.

NOT the graded submission (that stays exactly 2 slides per submission.pdf,
untouched, in build_ppt.py). This is the 13-slide explainer's content folded
into 5 pages for anyone who wants the full story without 13 pages of it -
still plain-language, still glossing every term, just tighter.

    python build_ppt_converged.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INK = RGBColor(0x1A, 0x22, 0x2C)
MUTED = RGBColor(0x5B, 0x66, 0x78)
FAINT = RGBColor(0x8A, 0x93, 0xA3)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
WARM = RGBColor(0xE8, 0x71, 0x22)
GOOD = RGBColor(0x1B, 0x7F, 0x5A)
BAD = RGBColor(0xB4, 0x33, 0x2E)
RULE = RGBColor(0xD8, 0xDE, 0xE6)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
NAVY = RGBColor(0x11, 0x24, 0x38)

W, H = Inches(13.333), Inches(7.5)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, bold=False, color=INK, space_before=0, space_after=2,
       first=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def _rect(slide, x, y, w, h, fill, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = RULE
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, kicker, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(0.09), ACCENT)
    tf = _box(slide, Inches(0.6), Inches(0.32), Inches(12.2), Inches(1.05))
    _p(tf, kicker, 11, bold=True, color=ACCENT, first=True)
    _p(tf, title, 25, bold=True, color=INK, space_before=1)
    if subtitle:
        _p(tf, subtitle, 12, color=MUTED, space_before=3)
    _rect(slide, Inches(0.6), Inches(1.36), Inches(12.13), Inches(0.012), RULE)


def footer(slide, n, total=5):
    tf = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf, "Visual Anomaly Detection — Converged Overview", 9.5, color=FAINT, first=True)
    tf2 = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf2, f"{n} / {total}", 9.5, color=FAINT, first=True, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ================================================================== SLIDE 1
def slide_title_problem(prs):
    s = new_slide(prs)
    _bg(s, NAVY)
    _rect(s, 0, Inches(6.9), W, Inches(0.6), ACCENT)

    tf = _box(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5))
    _p(tf, "VISUAL ANOMALY DETECTION — PLAIN-LANGUAGE OVERVIEW", 12.5, bold=True,
       color=ACCENT, first=True)
    tf2 = _box(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.3))
    _p(tf2, "Teaching a Camera to Notice What's Wrong", 32, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF), first=True)

    y = Inches(2.55)
    _rect(s, Inches(0.9), y, Inches(11.5), Inches(1.0), RGBColor(0x1B, 0x30, 0x4A))
    t = _box(s, Inches(1.2), y + Inches(0.16), Inches(11.0), Inches(0.7))
    _p(t, "THE IDEA IN ONE SENTENCE", 10, bold=True, color=ACCENT, first=True)
    _p(t, "Watch every video cheaply for most moments; call in an expensive expert "
          "only for the handful that actually look wrong.", 14, color=RGBColor(0xFF, 0xFF, 0xFF),
       space_before=3)

    y2 = Inches(3.85)
    cards = [
        ("The everyday version",
         "A parked car is normal. The SAME-LOOKING car stopped in a moving highway "
         "lane is not. The difference is context — where it is, and for how long — "
         "not what's in the photo."),
        ("The hardware version",
         "One ordinary gaming laptop, no data-centre GPU. A 4GB graphics card is a "
         "small kitchen counter — the biggest AI models simply don't fit on it."),
        ("So it was two problems at once",
         "Understand context well enough to know what's abnormal HERE, and do it "
         "cheaply enough to run constantly, on modest hardware, in real time."),
    ]
    cw = Inches(3.75)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.9) + i * (cw + Inches(0.12))
        _rect(s, x, y2, cw, Inches(2.55), RGBColor(0x16, 0x28, 0x40))
        t = _box(s, x + Inches(0.22), y2 + Inches(0.2), cw - Inches(0.44), Inches(2.2))
        _p(t, h, 13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=True)
        _p(t, b, 11, color=RGBColor(0xB8, 0xC4, 0xD6), space_before=6)


# ================================================================== SLIDE 2
def slide_architecture(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "THE ARCHITECTURE", "Four Checkpoints, Each Cheaper Than an Expert Opinion",
           "A triage line — a nurse checks everyone first; the doctor only sees who needs one.")

    y = Inches(1.58)
    stages = [
        ("1 · every frame", "\"What's here, and where?\"",
         "A small, fast detector draws a box around every vehicle and person and "
         "gives each one an ID it keeps across frames — retrained on overhead "
         "drone footage, since street-level training missed most objects seen "
         "from altitude."),
        ("2 · every frame, no AI", "\"Has it been sitting too long?\"",
         "Pure counting: how many seconds has this ID stayed put, how fast is it "
         "moving, which zone is it in. A stopwatch and a map — checkable facts, "
         "not a model's guess. This is what decides yes/no."),
        ("3 · a few sampled frames", "\"Fire, smoke, a flood?\"",
         "A small image model glances at a handful of frames for things a SINGLE "
         "photo genuinely can show. Deliberately narrow — it never has to judge "
         "motion or duration."),
        ("4 · rare, expensive calls", "\"Expert, describe this.\"",
         "A vision-language model — reads an image, writes a sentence — is "
         "called on well under 2% of frames: only what checkpoints 1-3 already "
         "flagged. It can raise an alarm, never silently clear one."),
    ]
    bw, gap = Inches(2.92), Inches(0.18)
    for i, (kicker, q, body) in enumerate(stages):
        x = Inches(0.6) + i * (bw + gap)
        _rect(s, x, y, bw, Inches(3.15), CARD, line=True)
        _rect(s, x, y, bw, Inches(0.06), WARM if i == 3 else ACCENT)
        t = _box(s, x + Inches(0.2), y + Inches(0.2), bw - Inches(0.4), Inches(2.8))
        _p(t, kicker.upper(), 9.5, bold=True, color=WARM if i == 3 else ACCENT, first=True)
        _p(t, q, 13, bold=True, color=INK, space_before=7, italic=True)
        _p(t, body, 10.5, color=MUTED, space_before=9)
        if i < 3:
            a = _box(s, x + bw + Inches(0.01), y + Inches(1.4), gap, Inches(0.3))
            _p(a, "›", 18, bold=True, color=RULE, first=True, align=PP_ALIGN.CENTER)

    y2 = Inches(4.98)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.9), RGBColor(0xFD, 0xF3, 0xE9), line=True)
    t = _box(s, Inches(0.95), y2 + Inches(0.18), Inches(11.4), Inches(1.55))
    _p(t, "WHY THIS ORDER, AND WHY WE CHOSE IT", 10.5, bold=True, color=WARM, first=True)
    _p(t, "Each checkpoint only handles what the ones before it couldn't rule out — by "
          "the time the expensive expert runs, three free filters already threw away "
          "almost everything unremarkable. It also follows one hard fact: a still "
          "photo cannot show motion. A moving car and a stopped one are pixel-"
          "identical in one frame. Asked directly, our expert model judged motion "
          "correctly only as often as a coin flip — so the yes/no decision moved to "
          "counting (checkpoint 2), which can measure duration with total reliability, "
          "and the AI's job narrowed to \"describe what you see\", which it's "
          "genuinely good at.", 11.5, color=INK, space_before=5)
    footer(s, 2)


# ================================================================== SLIDE 3
def slide_reasoning(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "THE REASONING, STEP BY STEP", "Every Design Choice Had a Concrete Reason",
           "Not \"best practice\" — a specific problem we hit, and the fix that followed.")

    rows = [
        ("Camera stabilisation",
         "A drone drifts in the air, so even a PARKED car appears to move across "
         "the picture.",
         "We measure how much the whole scene shifted between frames and cancel "
         "it out first — like image stabilisation, run in reverse."),
        ("Timing over snapshots",
         "A jam that clears in 5 seconds is normal; one that doesn't clear for 2 "
         "minutes is a problem — a difference no single frame can show.",
         "The counting stage keeps a running clock per tracked object, so "
         "\"how long\" becomes a number we compare to a threshold."),
        ("Our own retrained detector",
         "Off-the-shelf detectors, trained on eye-level photos, missed most "
         "objects seen from directly overhead.",
         "Retrained specifically on aerial footage. How much it actually finds "
         "nearly tripled, for the cost of one training run."),
        ("The expert AI can only escalate, never clear",
         "Trusting it both ways means one misread empty-looking frame could "
         "silently cancel an alert the counting stage already correctly raised.",
         "It can make an alert MORE serious on a visible danger, but the worst "
         "it can do otherwise is add an unneeded description — never erase a "
         "real catch."),
        ("We refused to tune against the one test set we could see",
         "It's tempting to nudge settings until they score perfectly on visible "
         "videos — but that number stops meaning anything on videos never seen.",
         "Every tuning idea was checked against held-out videos it hadn't "
         "touched. Several that looked great on visible data got measurably "
         "worse held-out, and were dropped."),
    ]
    y = Inches(1.56)
    rh = Inches(1.03)
    for i, (h, why, fix) in enumerate(rows):
        yy = y + rh * i
        _rect(s, Inches(0.6), yy, Inches(12.13), rh - Inches(0.06), CARD, line=True)
        num = _box(s, Inches(0.76), yy + Inches(0.1), Inches(0.5), Inches(0.5))
        _p(num, str(i + 1), 17, bold=True, color=ACCENT, first=True)
        t = _box(s, Inches(1.3), yy + Inches(0.07), Inches(11.35), rh - Inches(0.18))
        _p(t, h, 11.5, bold=True, color=INK, first=True)
        p = t.add_paragraph(); p.space_before = Pt(2)
        r = p.add_run(); r.text = "Problem: "
        r.font.size = Pt(9.6); r.font.bold = True; r.font.color.rgb = BAD; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = why
        r2.font.size = Pt(9.6); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"
        p2 = t.add_paragraph(); p2.space_before = Pt(1)
        r3 = p2.add_run(); r3.text = "Fix: "
        r3.font.size = Pt(9.6); r3.font.bold = True; r3.font.color.rgb = GOOD; r3.font.name = "Segoe UI"
        r4 = p2.add_run(); r4.text = fix
        r4.font.size = Pt(9.6); r4.font.color.rgb = MUTED; r4.font.name = "Segoe UI"
    footer(s, 3)


# ================================================================== SLIDE 4
def slide_results_tried(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "RESULTS & AN HONEST DEAD END", "What Works, What Doesn't, and One Idea We Killed",
           "Every number below is measured, not claimed.")

    y = Inches(1.56)
    _rect(s, Inches(0.6), y, Inches(5.95), Inches(2.65), CARD, line=True)
    t = _box(s, Inches(0.88), y + Inches(0.18), Inches(5.4), Inches(2.3))
    _p(t, "SPEED & ACCURACY", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "Fast enough to watch live video as it happens: ~15 frames a second on a "
          "modest laptop card, against the ~12-13 needed to look current.", 11.5,
       color=MUTED, first=False, space_before=6)
    _p(t, "Solid at WHICH category of problem occurred on short clips. Weaker at "
          "pinpointing the exact start/end moment inside long videos — that needs "
          "a kind of cross-time memory our fastest component doesn't have.",
       11.5, color=MUTED, space_before=8)

    _rect(s, Inches(6.78), y, Inches(5.95), Inches(2.65), CARD, line=True)
    t2 = _box(s, Inches(7.06), y + Inches(0.18), Inches(5.4), Inches(2.3))
    _p(t2, "ARENA SCOREBOARD", 10.5, bold=True, color=ACCENT, first=True)
    board = [("Clip classification", "70% correct", GOOD),
             ("When it happens (short)", "72% correct", GOOD),
             ("When it happens (long)", "20% correct", BAD)]
    for i, (a, b, col) in enumerate(board):
        yy = y + Inches(0.52) + Inches(0.4) * i
        tr = _box(s, Inches(7.06), yy, Inches(5.4), Inches(0.36))
        p = tr.paragraphs[0]
        r = p.add_run(); r.text = a
        r.font.size = Pt(12); r.font.color.rgb = INK; r.font.name = "Segoe UI"
        r = p.add_run(); r.text = "     " + b
        r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = col; r.font.name = "Segoe UI"
    _p(_box(s, Inches(7.06), y + Inches(1.9), Inches(5.4), Inches(0.6)),
       "We built a tool that recalculates our score exactly like the official "
       "scoreboard, before every submission — no guessing.", 10.5, color=MUTED, first=True)

    y2 = Inches(4.42)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(2.35), CARD, line=True)
    _rect(s, Inches(0.6), y2, Inches(0.06), Inches(2.35), BAD)
    t3 = _box(s, Inches(0.95), y2 + Inches(0.18), Inches(11.4), Inches(2.0))
    _p(t3, "WHAT WE TRIED AND DID NOT KEEP", 10.5, bold=True, color=BAD, first=True)
    _p(t3, "The idea: since a still frame can't show motion, feed the model the "
          "DIFFERENCE between frames instead — anything not moving cancels to black.",
       12, color=INK, space_before=6)
    _p(t3, "What we measured: on a 6-minute video where only 2% was truly anomalous, "
          "the new model flagged 100% of it — worse than before. In a scene full of "
          "ordinary moving traffic, almost everything is \"moving\" all the time, so "
          "the real clue (one person standing unusually still) never stood out any "
          "more than it did in a plain photo. We measured it honestly, saw it made "
          "things worse, and left it out.", 12, color=MUTED, space_before=6)
    footer(s, 4)


# ================================================================== SLIDE 5
def slide_limits_glossary(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "LIMITATIONS & QUICK GLOSSARY", "What's Left, and Every Term Used So Far",
           "Naming the gap plainly beats a vague claim that everything works.")

    y = Inches(1.56)
    _rect(s, Inches(0.6), y, Inches(6.9), Inches(3.35), CARD, line=True)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(3.35), BAD)
    t = _box(s, Inches(0.92), y + Inches(0.2), Inches(6.3), Inches(2.95))
    _p(t, "THE MAIN GAP", 10.5, bold=True, color=BAD, first=True)
    _p(t, "Pinpointing exactly when a long problem starts and ends.", 13.5, bold=True,
       color=INK, space_before=5)
    _p(t, "In a 6-minute video where a problem lasts 7 seconds, every fast checkpoint "
          "can only ask \"does this MOMENT look like the category?\" — and a busy "
          "scene often answers \"kind of, yes\" throughout, since busy scenes always "
          "have some motion. The counting stage (checkpoint 2) already keeps a "
          "per-person clock that knows the real answer; the unfinished work is "
          "letting that clock decide the exact timing, instead of a single-photo "
          "judgement that was never built to give one.", 11.5, color=MUTED, space_before=6)
    _p(t, "Two smaller, data-limited gaps: \"stalled vehicle\" had only 4 example "
          "videos to learn from, and \"traffic congestion\" had 23 — both too few "
          "for a model to learn reliably, regardless of design.", 11, color=MUTED, space_before=8)

    x2 = Inches(7.73)
    _rect(s, x2, y, Inches(5.0), Inches(3.35), CARD, line=True)
    t2 = _box(s, x2 + Inches(0.24), y + Inches(0.18), Inches(4.55), Inches(3.0))
    _p(t2, "GLOSSARY", 10.5, bold=True, color=ACCENT, first=True)
    terms = [
        ("Model", "a program that learned a task from examples."),
        ("Detection / tracking", "finding an object, then keeping the same ID for it over time."),
        ("Vision-language model", "an AI that looks at an image and writes a sentence about it."),
        ("Fine-tuning", "retraining an existing model on examples specific to one job."),
        ("GPU / video memory", "the chip AI models run on, and its own separate memory."),
        ("False alarm", "flagging a problem that wasn't real — the costlier mistake here."),
        ("Recall", "of all the real problems, what fraction did the system find?"),
        ("Ground truth", "the verified correct answer, used to check the system."),
    ]
    for term, defn in terms:
        p = t2.add_paragraph(); p.space_before = Pt(4)
        r = p.add_run(); r.text = term + " — "
        r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = defn
        r2.font.size = Pt(10.5); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"

    y2 = y + Inches(3.5)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.05), RGBColor(0xEE, 0xF4, 0xF8), line=True)
    t3 = _box(s, Inches(0.92), y2 + Inches(0.14), Inches(11.5), Inches(0.8))
    _p(t3, "WITH MORE TIME: ", 11, bold=True, color=ACCENT, first=True)
    p = t3.paragraphs[0]
    r = p.add_run()
    r.text = ("give the per-person clock the final say on WHEN an event happened, and "
              "reserve the quick-glance model for WHAT category it belongs to — the job "
              "each one is actually suited for.")
    r.font.size = Pt(11); r.font.color.rgb = INK; r.font.name = "Segoe UI"
    footer(s, 5)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_title_problem(prs)
    slide_architecture(prs)
    slide_reasoning(prs)
    slide_results_tried(prs)
    slide_limits_glossary(prs)
    out = Path("deck")
    out.mkdir(exist_ok=True)
    f = out / "FlytBase_AHC_Overview_5pg.pptx"
    try:
        prs.save(f)
    except PermissionError:
        f = out / "FlytBase_AHC_Overview_5pg_v2.pptx"
        prs.save(f)
    print(f"wrote {f}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
