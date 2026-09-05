"""Generate the ELABORATIVE explainer deck — 5 slides, plain language.

NOT the graded submission. submission.pdf requires the uploaded PPT to be
exactly 2 slides ("Keep the final presentation to 2 slides only") - that file
is build_ppt.py's output and is untouched by this script. This is a separate,
longer walkthrough for anyone (technical or not) who wants the reasoning
behind every step explained without jargon - a teammate, a judge who wants
depth after the 2-slide pitch, a reviewer.

Register: every technical term gets a plain-English gloss inline, the way you
would actually explain it to someone out loud. No acronym appears without its
meaning attached the first time.

    python build_ppt_explainer.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INK = RGBColor(0x1A, 0x22, 0x2C)
MUTED = RGBColor(0x5B, 0x66, 0x78)
FAINT = RGBColor(0x8A, 0x93, 0xA3)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)     # calmer blue — explainer register, not pitch-deck amber
WARM = RGBColor(0xE8, 0x71, 0x22)
GOOD = RGBColor(0x1B, 0x7F, 0x5A)
BAD = RGBColor(0xB4, 0x33, 0x2E)
RULE = RGBColor(0xD8, 0xDE, 0xE6)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
NAVY = RGBColor(0x11, 0x24, 0x38)

W, H = Inches(13.333), Inches(7.5)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, bold=False, color=INK, space_before=0, space_after=2,
       first=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def _rect(slide, x, y, w, h, fill, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = RULE
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, kicker, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(0.09), ACCENT)
    tf = _box(slide, Inches(0.6), Inches(0.34), Inches(12.2), Inches(1.05))
    _p(tf, kicker, 11.5, bold=True, color=ACCENT, first=True)
    _p(tf, title, 26, bold=True, color=INK, space_before=1)
    if subtitle:
        _p(tf, subtitle, 12.5, color=MUTED, space_before=3)
    _rect(slide, Inches(0.6), Inches(1.4), Inches(12.13), Inches(0.012), RULE)


def footer(slide, n, total=13):
    tf = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf, "Visual Anomaly Detection — Explainer", 9.5, color=FAINT, first=True)
    tf2 = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf2, f"{n} / {total}", 9.5, color=FAINT, first=True, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ================================================================== SLIDE 1
def slide_title(prs):
    s = new_slide(prs)
    _bg(s, NAVY)
    _rect(s, 0, Inches(6.9), W, Inches(0.6), ACCENT)

    tf = _box(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5))
    _p(tf, "AN EXPLAINER, IN PLAIN LANGUAGE", 13, bold=True, color=ACCENT, first=True)

    tf2 = _box(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.9))
    _p(tf2, "Visual Anomaly Detection", 42, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=True)
    _p(tf2, "Teaching a Camera to Notice What's Wrong", 30, bold=True,
       color=RGBColor(0xC9, 0xD2, 0xE0), space_before=4)

    tf3 = _box(s, Inches(0.9), Inches(4.15), Inches(10.8), Inches(1.3))
    _p(tf3, "This deck answers three questions in order: what problem are we solving, "
            "what did we build to solve it, and why did we build it that way instead "
            "of some other way — explained so that no prior AI background is needed.",
       15, color=RGBColor(0xC9, 0xD2, 0xE0), first=True, italic=True)

    # the one-sentence idea, as a lead-in
    y = Inches(5.7)
    _rect(s, Inches(0.9), y, Inches(11.5), Inches(1.0), RGBColor(0x1B, 0x30, 0x4A))
    t = _box(s, Inches(1.2), y + Inches(0.18), Inches(11.0), Inches(0.7))
    _p(t, "THE IDEA IN ONE SENTENCE", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "Watch every video cheaply for the first ninety seconds of a task; call in "
          "an expensive expert only for the handful of moments that actually look wrong.",
       14.5, color=RGBColor(0xFF, 0xFF, 0xFF), space_before=4)


# ================================================================== SLIDE 2
def slide_problem(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "THE PROBLEM, IN PLAIN TERMS", "What Are We Actually Asking a Computer to Do?",
           "Not \"is there a car in this picture\" — a much harder question about context.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(6.0), Inches(2.55), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.22), Inches(5.4), Inches(2.15))
    _p(t, "THE EVERYDAY VERSION", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "A parked car is normal. A car stopped in the middle of a moving "
          "highway lane is not.", 14, bold=True, color=INK, space_before=6)
    _p(t, "The two situations can look almost identical in a single photo — a car, "
          "sitting still. What makes one an emergency and the other unremarkable is "
          "context: where it is, and for how long it has been there. A computer "
          "has to learn to ask those questions, not just \"do I see a car?\"",
       12, color=MUTED, space_before=8)

    _rect(s, Inches(6.83), y, Inches(5.9), Inches(2.55), CARD, line=True)
    t2 = _box(s, Inches(7.13), y + Inches(0.22), Inches(5.3), Inches(2.15))
    _p(t2, "THE HARDWARE VERSION", 10.5, bold=True, color=ACCENT, first=True)
    _p(t2, "We had one ordinary gaming laptop — no data-centre GPU.", 14, bold=True,
       color=INK, space_before=6)
    _p(t2, "A GTX 1650 with 4 GB of video memory — think of it as a small "
          "kitchen counter to work on, not an industrial one. Powerful AI vision "
          "models don't fit on a counter that size, and running one on every "
          "single frame of video would be far too slow to matter in real time.",
       12, color=MUTED, space_before=8)

    y2 = Inches(4.5)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(2.25), RGBColor(0xEE, 0xF4, 0xF8), line=True)
    _rect(s, Inches(0.6), y2, Inches(0.06), Inches(2.25), ACCENT)
    t3 = _box(s, Inches(0.95), y2 + Inches(0.22), Inches(11.4), Inches(1.85))
    _p(t3, "SO THE REAL TASK WAS TWO PROBLEMS AT ONCE", 10.5, bold=True, color=ACCENT, first=True)
    pairs = [
        ("Understanding context", "teach the system what \"normal\" looks like for a "
         "given place, so it can spot what breaks the pattern — not just detect objects."),
        ("Working within a budget", "do this on modest hardware, fast enough to matter, "
         "without needing a warehouse full of expensive computers."),
    ]
    for i, (h, b) in enumerate(pairs):
        yy = y2 + Inches(0.62) + Inches(0.66) * i
        tr = _box(s, Inches(0.95), yy, Inches(11.0), Inches(0.6))
        p = tr.paragraphs[0]
        r = p.add_run(); r.text = f"{i+1}.  {h} — "
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = b
        r2.font.size = Pt(13); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"
    footer(s, 2)


# ================================================================== SLIDE 3
def slide_hard(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "WHY THIS IS GENUINELY HARD", "Three Reasons This Isn't Just \"Run an App Store AI Model\"",
           "Each one ruled out the obvious, simplest answer.")

    cards = [
        ("\"Normal\" depends on WHERE you are",
         "A crowd of people standing still is completely normal outside a stadium "
         "gate, and a clear sign of trouble in the middle of a road. The same "
         "picture means two different things depending on the location — so a "
         "model can't just learn \"what does a problem look like\" in general."),
        ("Most problems are about TIME, not appearance",
         "A car in a lane is not the story — a car that HASN'T MOVED for two "
         "minutes is. That fact doesn't exist in any single photo; it only "
         "exists when you compare many moments to each other."),
        ("Powerful AI models are too slow and too heavy to run constantly",
         "The best available vision-language models (models that can look at an "
         "image and describe it in a sentence) take tens of seconds per look and "
         "need more memory than our laptop's graphics card has. Running one on "
         "every frame of every video was never on the table."),
    ]
    y = Inches(1.72)
    cw = Inches(3.9)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.6) + i * (cw + Inches(0.2))
        _rect(s, x, y, cw, Inches(3.4), CARD, line=True)
        _rect(s, x, y, cw, Inches(0.06), ACCENT)
        t = _box(s, x + Inches(0.24), y + Inches(0.24), cw - Inches(0.48), Inches(3.0))
        _p(t, h, 14.5, bold=True, color=INK, first=True)
        _p(t, b, 12, color=MUTED, space_before=10)

    y2 = Inches(5.35)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.5), RGBColor(0xEE, 0xF4, 0xF8), line=True)
    t = _box(s, Inches(0.95), y2 + Inches(0.2), Inches(11.4), Inches(1.15))
    _p(t, "PUT TOGETHER", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "We needed a system that understands context, reasons about time, and stays "
          "cheap enough to run constantly — three requirements that pull against each "
          "other, which is exactly why a single AI model can't do this job alone.",
       13, color=INK, space_before=6)
    footer(s, 3)


# ================================================================== SLIDE 4
def slide_architecture(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "THE ARCHITECTURE", "Four Checkpoints, Each Cheaper Than an Expert Opinion",
           "Think of it as a hospital triage line — a nurse checks everyone first; the doctor only sees who needs one.")

    y = Inches(1.72)
    stages = [
        ("CHECKPOINT 1 · every single frame", "\"What's in the picture, and where?\"",
         "A fast, lightweight detector (a smaller AI model, tuned for speed) draws a "
         "box around every vehicle and person and gives each one an ID number it keeps "
         "across frames — so we know it's the SAME car two seconds later, not a new one."),
        ("CHECKPOINT 2 · every frame, pure arithmetic", "\"Has anything been sitting still too long?\"",
         "No AI model runs here at all — just counting. How many seconds has this "
         "car ID stayed in one place? Is it in a driving lane or a parking spot? "
         "That's a stopwatch and a map, not a neural network."),
        ("CHECKPOINT 3 · a handful of sampled frames", "\"Does this look like fire, smoke, "
         "or a flood?\"",
         "A small image-recognition model — trained specifically on these categories — "
         "glances at a few frames. It only handles things a SINGLE photo can show; it "
         "never has to judge motion or duration."),
        ("CHECKPOINT 4 · rare, expensive calls only", "\"An expert, please describe "
         "what's actually happening.\"",
         "A vision-language model (an AI that can look at an image and write a "
         "sentence about it) is called on well under 2% of frames — only the ones "
         "the earlier, cheap checkpoints flagged as suspicious."),
    ]
    bw, gap = Inches(2.92), Inches(0.18)
    for i, (kicker, q, body) in enumerate(stages):
        x = Inches(0.6) + i * (bw + gap)
        _rect(s, x, y, bw, Inches(3.55), CARD, line=True)
        _rect(s, x, y, bw, Inches(0.06), WARM if i == 3 else ACCENT)
        t = _box(s, x + Inches(0.2), y + Inches(0.22), bw - Inches(0.4), Inches(3.1))
        _p(t, kicker, 9.5, bold=True, color=WARM if i == 3 else ACCENT, first=True)
        _p(t, q, 13, bold=True, color=INK, space_before=8, italic=True)
        _p(t, body, 10.8, color=MUTED, space_before=10)
        if i < 3:
            a = _box(s, x + bw + Inches(0.01), y + Inches(1.6), gap, Inches(0.3))
            _p(a, "›", 18, bold=True, color=RULE, first=True, align=PP_ALIGN.CENTER)

    y2 = Inches(5.55)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.3), RGBColor(0xFD, 0xF3, 0xE9), line=True)
    t = _box(s, Inches(0.95), y2 + Inches(0.18), Inches(11.4), Inches(1.0))
    _p(t, "WHY IN THIS ORDER", 10.5, bold=True, color=WARM, first=True)
    _p(t, "Each checkpoint only has to handle what the ones before it couldn't rule out. "
          "By the time the expensive expert (checkpoint 4) is asked anything, three "
          "cheap filters have already thrown away almost everything unremarkable — "
          "which is the only reason this runs fast enough to matter at all.",
       12.5, color=INK, space_before=4)
    footer(s, 4)


# ================================================================== SLIDE 5
def slide_stage1_deep(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "CHECKPOINT 1, IN DEPTH", "Finding Things and Keeping Track of Them",
           "Two separate jobs bundled into one step: spotting objects, and remembering which object is which.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(5.9), Inches(3.4), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.22), Inches(5.3), Inches(3.0))
    _p(t, "SPOTTING (\"DETECTION\")", 11, bold=True, color=ACCENT, first=True)
    _p(t, "A small, fast AI model draws a box around every car, motorbike and "
          "person it can find, in every single frame.", 13, bold=True, color=INK, space_before=6)
    _p(t, "We deliberately used a SMALL model here, not the biggest one available. "
          "It has to run dozens of times per second, so speed matters more than "
          "squeezing out the last bit of accuracy — this is the \"nurse\" in the "
          "triage line, not the specialist.", 12, color=MUTED, space_before=8)
    _p(t, "We also retrained it on footage shot from directly overhead, because "
          "the version trained on ordinary street-level photos missed most objects "
          "seen from a drone's altitude — everything looks small and unfamiliar "
          "from straight above.", 12, color=MUTED, space_before=8)

    _rect(s, Inches(6.63), y, Inches(6.1), Inches(3.4), CARD, line=True)
    t2 = _box(s, Inches(6.93), y + Inches(0.22), Inches(5.5), Inches(3.0))
    _p(t2, "REMEMBERING (\"TRACKING\")", 11, bold=True, color=ACCENT, first=True)
    _p(t2, "Every box gets an ID number that sticks with it across frames, so "
          "\"car #7\" is the same car two seconds from now.", 13, bold=True, color=INK, space_before=6)
    _p(t2, "Without this, the system has no memory — it would see a car, then "
          "see A car a moment later, with no way to know if it's the same one "
          "standing still or a different one passing through. Everything the "
          "system later says about \"how long has this been stationary\" depends "
          "entirely on this step working.", 12, color=MUTED, space_before=8)
    _p(t2, "One extra correction: a drone drifts in the air, so the picture itself "
          "shifts slightly frame to frame. We measure and cancel out that camera "
          "movement first — otherwise a parked car would look like it's slowly "
          "sliding across the screen.", 12, color=MUTED, space_before=8)
    footer(s, 5)


# ================================================================== SLIDE 6
def slide_stage2_deep(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "CHECKPOINTS 2 & 3, IN DEPTH", "Counting, Then a Quick Glance",
           "No AI at all for the boolean decision — then a small, focused specialist for what a photo genuinely can show.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(5.9), Inches(3.4), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.22), Inches(5.3), Inches(3.0))
    _p(t, "CHECKPOINT 2 · COUNTING", 11, bold=True, color=ACCENT, first=True)
    _p(t, "For every tracked object: how many seconds has it stayed put? How "
          "fast is it moving? Which zone is it in?", 13, bold=True, color=INK, space_before=6)
    _p(t, "This is arithmetic, not artificial intelligence — a stopwatch and a "
          "map. \"This car has not moved in 12 seconds, and it's in a driving "
          "lane, not a parking spot\" is a fact you can check with a ruler, not "
          "a guess a model makes. That's exactly why it's trustworthy enough to "
          "be given the final yes/no decision.", 12, color=MUTED, space_before=8)
    _p(t, "This step also builds the traffic-jam, wrong-direction, and loitering "
          "(lingering-too-long) checks — all of them are really just \"how long\" "
          "and \"how fast\" questions in disguise.", 12, color=MUTED, space_before=8)

    _rect(s, Inches(6.63), y, Inches(6.1), Inches(3.4), CARD, line=True)
    t2 = _box(s, Inches(6.93), y + Inches(0.22), Inches(5.5), Inches(3.0))
    _p(t2, "CHECKPOINT 3 · A QUICK GLANCE", 11, bold=True, color=ACCENT, first=True)
    _p(t2, "A small, specially-trained image-recognition model looks at a "
          "handful of frames for things a SINGLE photo genuinely can show.",
       13, bold=True, color=INK, space_before=6)
    _p(t2, "Fire, smoke, a flood, or debris on the road don't need motion or "
          "timing to recognise — a photograph is enough, the way you'd "
          "recognise a photo of a house fire instantly. So this step is kept "
          "deliberately narrow: it only ever answers questions a still image "
          "can honestly answer.", 12, color=MUTED, space_before=8)
    _p(t2, "It is trained specifically on these categories, on footage similar "
          "to what we expect to see — not a generic \"recognise anything\" model, "
          "which would be slower and less accurate at this one job.",
       12, color=MUTED, space_before=8)
    footer(s, 6)


# ================================================================== SLIDE 7
def slide_stage3_deep(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "CHECKPOINT 4, IN DEPTH", "Calling in the Expert — Rarely, and on Purpose",
           "The one step that uses a big, expensive AI model — deliberately limited to under 2% of frames.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(12.13), Inches(1.7), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.2), Inches(11.5), Inches(1.35))
    _p(t, "WHAT IT DOES", 10.5, bold=True, color=WARM, first=True)
    _p(t, "A vision-language model — an AI that looks at a picture and writes a "
          "sentence describing it — is shown the specific frame that Checkpoints "
          "1-3 flagged as suspicious, along with a short note on what triggered "
          "the flag. It writes back a description, and it can raise the alarm "
          "level if it spots something dangerous the earlier steps missed.",
       12.5, color=INK, space_before=4)

    y2 = Inches(3.62)
    cards = [
        ("Why so rarely?",
         "This model is slow (tens of seconds per look) and needs more computer "
         "memory than our laptop's graphics card has to spare alongside "
         "Checkpoint 1's detector. Calling it on every frame was never fast "
         "enough to be useful — so it's saved for the rare frames that need it."),
        ("Why can it raise alarms but never lower them?",
         "It can describe a scene well, but — as the next slide explains in "
         "detail — a single picture genuinely cannot tell it whether an object "
         "has been stopped for two seconds or two minutes. So it is trusted to "
         "spot dangers, not to overrule a timing measurement it cannot see."),
        ("Was it customised for this job?",
         "Yes — it was specifically retrained (a process called fine-tuning) on "
         "examples from this exact kind of footage, so its descriptions and "
         "judgement are sharper on aerial traffic scenes than an untrained, "
         "off-the-shelf version would be."),
    ]
    cw = Inches(3.9)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.6) + i * (cw + Inches(0.2))
        _rect(s, x, y2, cw, Inches(2.9), CARD, line=True)
        t2 = _box(s, x + Inches(0.24), y2 + Inches(0.2), cw - Inches(0.48), Inches(2.55))
        _p(t2, h, 13, bold=True, color=INK, first=True)
        _p(t2, b, 11.5, color=MUTED, space_before=8)
    footer(s, 7)


# ================================================================== SLIDE 8
def slide_why(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "WHY WE CHOSE THIS APPROACH", "The One Fact That Decided Almost Everything",
           "A single frozen photo cannot show you motion — and most of what we needed to detect IS motion.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(12.13), Inches(2.15), CARD, line=True)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(2.15), WARM)
    t = _box(s, Inches(0.95), y + Inches(0.22), Inches(11.4), Inches(1.75))
    _p(t, "A PHOTO OF A STOPPED CAR AND A PHOTO OF A MOVING CAR LOOK IDENTICAL",
       12, bold=True, color=WARM, first=True)
    _p(t, "We tried simply asking the expert vision-language model, directly: "
          "\"is this anomalous?\" It got the answer right only as often as a coin "
          "flip would — because from one still image, it genuinely cannot tell "
          "whether a car is parked or paused at a red light. Its written "
          "description was accurate (\"a car sitting in a driving lane while "
          "traffic flows around it\"), but its yes/no answer contradicted its own "
          "sentence, because a single picture just doesn't contain that "
          "information.", 12.5, color=MUTED, space_before=6)

    y2 = Inches(4.1)
    cards = [
        ("So the yes/no decision moved to counting, not to AI",
         "Whether something is anomalous is decided by simple measurements — how "
         "long has it stood still, which zone is it in, are the vehicles around it "
         "still moving — which a computer can measure with total reliability. "
         "Nothing here is guessed."),
        ("The AI's job became \"describe\", not \"judge\"",
         "The expert model is asked to say what it sees, in words — smoke, a "
         "person down, a crowd — a task it is genuinely good at. It can raise an "
         "alarm louder if it spots something dangerous, but it is never allowed "
         "to overrule a measurement that already found a problem."),
    ]
    cw = Inches(5.95)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.6) + i * (cw + Inches(0.23))
        _rect(s, x, y2, cw, Inches(2.65), CARD, line=True)
        tc = _box(s, x + Inches(0.28), y2 + Inches(0.22), cw - Inches(0.56), Inches(2.2))
        _p(tc, h, 14, bold=True, color=GOOD, first=True)
        _p(tc, b, 12.5, color=MUTED, space_before=8)
    footer(s, 8)


# ================================================================== SLIDE 9
def slide_reasoning(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "THE REASONING, STEP BY STEP", "Every Design Choice Had a Concrete Reason",
           "Not \"best practice\" — a specific problem we hit, and the fix that followed.")

    rows = [
        ("The camera stabilisation step",
         "A drone drifts and turns in the air, so even a PARKED car appears to move "
         "across the picture. Without correcting for the drone's own motion first, "
         "the system would think every car was driving, all the time.",
         "We calculate how much the whole scene shifted between frames and subtract "
         "that out — the same trick as image stabilisation in a phone camera, run in "
         "reverse to reveal what's truly still."),
        ("Why timing (\"has this been happening for a while?\") beats a snapshot",
         "A traffic jam that clears in five seconds is normal; one that doesn't "
         "clear for two minutes is a problem. The difference is entirely about "
         "duration, which a single frame cannot capture at all.",
         "The counting stage keeps a running clock on every tracked object, so "
         "\"how long\" becomes a number we can simply compare to a threshold."),
        ("Why we trained our own small object-detector instead of a general one",
         "Off-the-shelf detectors are trained mostly on eye-level photos — a "
         "person walking toward the camera, a car seen from the side. From directly "
         "above, at altitude, everything looks small and unfamiliar, and the "
         "off-the-shelf detector missed the majority of real objects.",
         "We retrained it specifically on aerial, top-down footage. Recall (how much "
         "it actually finds) nearly tripled, for the cost of a training run, not new "
         "hardware."),
        ("Why we let a description-writing AI \"escalate\" but never \"clear\" an alert",
         "If we trusted it both ways, a model that misreads an empty-looking "
         "frame as \"nothing here\" could silently cancel an alert that the "
         "counting stage had already correctly raised — turning our safety net "
         "into a hole in it.",
         "It is allowed to make an alert MORE serious if it spots a visible "
         "danger (fire, a person down), but it can never make one go away. The "
         "worst it can do is add an unnecessary description; it can't undo a "
         "real catch."),
        ("Why we refused to fine-tune our settings against the one test set we could see",
         "It's tempting to nudge every setting until it scores perfectly on the "
         "videos you have in front of you. But that number stops meaning anything "
         "the moment you test on videos you haven't seen — including the ones "
         "used for real judging.",
         "We checked every tuning idea against a held-out slice it had never "
         "touched before trusting it. Several ideas that looked great on the "
         "visible videos got measurably worse on the held-out ones, and we "
         "dropped them rather than keep the flattering number."),
    ]
    y = Inches(1.56)
    rh = Inches(1.1)
    for i, (h, why, fix) in enumerate(rows):
        yy = y + rh * i
        _rect(s, Inches(0.6), yy, Inches(12.13), rh - Inches(0.06), CARD, line=True)
        num = _box(s, Inches(0.76), yy + Inches(0.1), Inches(0.5), Inches(0.5))
        _p(num, str(i + 1), 18, bold=True, color=ACCENT, first=True)
        t = _box(s, Inches(1.32), yy + Inches(0.08), Inches(11.35), rh - Inches(0.2))
        _p(t, h, 12, bold=True, color=INK, first=True)
        p = t.add_paragraph(); p.space_before = Pt(2)
        r = p.add_run(); r.text = "The problem: "
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = BAD; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = why
        r2.font.size = Pt(10); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"
        p2 = t.add_paragraph(); p2.space_before = Pt(1)
        r3 = p2.add_run(); r3.text = "The fix: "
        r3.font.size = Pt(10); r3.font.bold = True; r3.font.color.rgb = GOOD; r3.font.name = "Segoe UI"
        r4 = p2.add_run(); r4.text = fix
        r4.font.size = Pt(10); r4.font.color.rgb = MUTED; r4.font.name = "Segoe UI"
    footer(s, 9)


# ================================================================== SLIDE 10
def slide_tried_rejected(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "WHAT WE TRIED AND DIDN'T KEEP", "An Honest Look at a Dead End",
           "Worth showing precisely because it was checked carefully, not just assumed to work.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(12.13), Inches(2.0), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.2), Inches(11.5), Inches(1.6))
    _p(t, "THE IDEA", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "A single video frame can't show motion — so what if we fed the model the "
          "DIFFERENCE between two frames instead of a plain photo? In theory, anything "
          "that isn't moving would cancel out to black, and only genuine motion would "
          "remain — which sounded like exactly the missing ingredient.",
       12.5, color=INK, space_before=6)

    y2 = Inches(3.9)
    cards = [
        ("What we expected",
         "A model trained on \"difference images\" would finally be able to tell "
         "the difference between a car that's driving through a busy scene and "
         "one that's been sitting still in the same busy scene."),
        ("What we actually measured",
         "On a six-minute test video where only 2% of the footage was truly "
         "anomalous, the new model flagged 100% of it as suspicious — worse "
         "than before we made the change, not better."),
        ("Why it failed, in plain terms",
         "In a scene full of ordinary moving traffic, almost everything is "
         "\"moving\" all the time. One person standing still in a busy scene "
         "doesn't stand out in the difference image any more than they stood "
         "out in the original photo — the real clue was never in one frame at "
         "all, difference or not."),
    ]
    cw = Inches(3.9)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.6) + i * (cw + Inches(0.2))
        _rect(s, x, y2, cw, Inches(2.5), CARD, line=True)
        _rect(s, x, y2, cw, Inches(0.06), BAD if i == 1 else ACCENT)
        t2 = _box(s, x + Inches(0.24), y2 + Inches(0.24), cw - Inches(0.48), Inches(2.1))
        _p(t2, h, 13, bold=True, color=INK, first=True)
        _p(t2, b, 11.5, color=MUTED, space_before=8)
    footer(s, 10)


# ================================================================== SLIDE 11
def slide_results(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "RESULTS, IN PLAIN LANGUAGE", "How Well Does It Actually Work?",
           "Numbers translated into what they mean, not just reported.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(5.95), Inches(3.0), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.22), Inches(5.4), Inches(2.6))
    _p(t, "SPEED", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "Fast enough to watch live video as it happens.", 14, bold=True, color=INK, space_before=6)
    _p(t, "The full system processes about 15 video frames every second on our "
          "modest laptop graphics card. Ordinary drone footage needs about 12-13 "
          "frames a second to look smooth and current — so we're comfortably "
          "ahead of \"live\", with room to spare.", 12, color=MUTED, space_before=8)
    _p(t, "The expert AI model (Checkpoint 4) is much slower on its own — tens "
          "of seconds per look — which is exactly why it's only used sparingly, "
          "not on every frame.", 12, color=MUTED, space_before=8)

    _rect(s, Inches(6.78), y, Inches(5.95), Inches(3.0), CARD, line=True)
    t2 = _box(s, Inches(7.08), y + Inches(0.22), Inches(5.4), Inches(2.6))
    _p(t2, "ACCURACY", 10.5, bold=True, color=ACCENT, first=True)
    _p(t2, "Strong at \"what happened\", weaker at \"exactly when\".", 14, bold=True,
       color=INK, space_before=6)
    _p(t2, "On short clips, and on correctly identifying WHICH category of "
          "problem occurred, the system performs solidly — a majority of these "
          "are answered correctly.", 12, color=MUTED, space_before=8)
    _p(t2, "It performs noticeably worse at pinpointing the exact start and end "
          "moment of an event inside long videos (many minutes), because that "
          "needs a kind of memory across time that our current fastest "
          "component doesn't have — explained on the next slide.",
       12, color=MUTED, space_before=8)

    y2 = Inches(5.0)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.85), RGBColor(0xEE, 0xF4, 0xF8), line=True)
    _rect(s, Inches(0.6), y2, Inches(0.06), Inches(1.85), ACCENT)
    t3 = _box(s, Inches(0.95), y2 + Inches(0.2), Inches(11.4), Inches(1.5))
    _p(t3, "WHY WE CHECK OUR OWN SCORE INSTEAD OF GUESSING", 10.5, bold=True, color=ACCENT, first=True)
    _p(t3, "We built a small tool that recalculates our score the same way the official "
          "scoreboard does, before we submit anything — so every change is measured "
          "against a real number beforehand, instead of \"we think this should help\" "
          "and finding out only after using up a submission.",
       12.5, color=INK, space_before=6)
    footer(s, 11)


# ================================================================== SLIDE 12
def slide_limits(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "LIMITATIONS, HONESTLY STATED", "What Doesn't Work Yet, and Why",
           "Naming the gap plainly is worth more than a vague claim that everything works.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(7.0), Inches(4.75), CARD, line=True)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(4.75), BAD)
    t = _box(s, Inches(0.95), y + Inches(0.24), Inches(6.4), Inches(4.3))
    _p(t, "THE MAIN GAP", 11, bold=True, color=BAD, first=True)
    _p(t, "Pinpointing exactly when a long-running problem starts and ends.", 15, bold=True,
       color=INK, space_before=6)
    _p(t, "Picture a six-minute video where someone is loitering suspiciously for "
          "just seven of those seconds. Every one of our fast checkpoints can only "
          "look at individual moments and ask \"does this look like the category "
          "loitering?\" — and the honest answer for a busy, ordinary scene is often "
          "\"kind of, yes\", the whole way through, because busy scenes always have "
          "SOME people moving around in them.",
       12.5, color=MUTED, space_before=8)
    _p(t, "The one component that actually knows \"this exact person has been "
          "standing here for a suspiciously long time\" is the counting stage "
          "(Checkpoint 2) — because it keeps a per-person clock. The unfinished "
          "work is connecting that per-person clock directly to the final "
          "start/end answer, instead of relying on the quick-glance model's "
          "single-photo judgement for timing it was never built to give.",
       12.5, color=MUTED, space_before=10)

    x2 = Inches(7.86)
    _rect(s, x2, y, Inches(4.87), Inches(2.85), CARD, line=True)
    t2 = _box(s, x2 + Inches(0.28), y + Inches(0.22), Inches(4.35), Inches(2.5))
    _p(t2, "TWO SMALLER, DATA-LIMITED GAPS", 11, bold=True, color=ACCENT, first=True)
    _p(t2, "A model can only learn a pattern it has SEEN many times.", 12.5, bold=True,
       color=INK, space_before=6)
    for item in ["\"Stalled or broken-down vehicle\" had only 4 example videos "
                 "in our entire training set — nowhere near enough to teach a "
                 "model reliably.",
                 "\"Traffic congestion\" had 23 — better, but still thin "
                 "compared to categories with over a hundred examples."]:
        p = t2.add_paragraph(); p.space_before = Pt(6)
        r = p.add_run(); r.text = "→ "
        r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"

    y2 = y + Inches(3.05)
    _rect(s, x2, y2, Inches(4.87), Inches(1.7), RGBColor(0xF0, 0xF3, 0xF7), line=True)
    t3 = _box(s, x2 + Inches(0.28), y2 + Inches(0.2), Inches(4.35), Inches(1.35))
    _p(t3, "WITH MORE TIME", 10, bold=True, color=MUTED, first=True)
    _p(t3, "Give the per-person clock the final say on WHEN an event happened, "
          "and reserve the quick-glance model for WHAT category it belongs to — "
          "the job each one is actually suited for.", 11.5, color=INK, space_before=6)
    footer(s, 12)


# ================================================================== SLIDE 13
def slide_glossary(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "A QUICK GLOSSARY", "Every Technical Term Used in This Deck",
           "For anyone reading this without an AI background — no term is assumed.")

    terms = [
        ("Model", "A program that has learned a task from examples, rather than "
         "being told explicit step-by-step rules."),
        ("Detection", "Finding and drawing a box around an object in an image — "
         "\"there is a car here.\""),
        ("Tracking", "Giving each detected object an ID number that stays the "
         "same across frames, so you know it's the same car, not a new one."),
        ("Vision-language model", "An AI that can look at an image and produce "
         "a written sentence describing it, or answer questions about it."),
        ("Fine-tuning", "Taking an existing, general AI model and retraining it "
         "further on examples specific to one job, to sharpen it for that job."),
        ("Inference", "The act of a trained model actually being run to produce "
         "an answer — as opposed to the earlier process of training it."),
        ("GPU / video memory", "The specialised chip (and its own separate "
         "memory) that AI models run on, because it's far faster at this kind "
         "of math than the everyday processor in a computer."),
        ("False alarm", "The system flags something as a problem when nothing "
         "was actually wrong — the more expensive kind of mistake to make here."),
        ("Recall", "Out of all the real problems that existed, what fraction did "
         "the system actually find? A low recall means it misses things."),
        ("Ground truth", "The verified, correct answer for a video — used to "
         "check whether the system's answer was right."),
    ]
    y = Inches(1.7)
    col_w = Inches(5.95)
    row_h = Inches(1.0)
    for i, (term, defn) in enumerate(terms):
        col = i // 5
        row = i % 5
        x = Inches(0.6) + col * (col_w + Inches(0.23))
        yy = y + row * row_h
        _rect(s, x, yy, col_w, row_h - Inches(0.1), CARD, line=True)
        t = _box(s, x + Inches(0.2), yy + Inches(0.12), col_w - Inches(0.4), row_h - Inches(0.3))
        _p(t, term, 12.5, bold=True, color=ACCENT, first=True)
        _p(t, defn, 10.8, color=MUTED, space_before=2)
    footer(s, 13)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_title(prs)
    slide_problem(prs)
    slide_hard(prs)
    slide_architecture(prs)
    slide_stage1_deep(prs)
    slide_stage2_deep(prs)
    slide_stage3_deep(prs)
    slide_why(prs)
    slide_reasoning(prs)
    slide_tried_rejected(prs)
    slide_results(prs)
    slide_limits(prs)
    slide_glossary(prs)
    out = Path("deck")
    out.mkdir(exist_ok=True)
    f = out / "FlytBase_AHC_Explainer.pptx"
    try:
        prs.save(f)
    except PermissionError:
        f = out / "FlytBase_AHC_Explainer_v2.pptx"
        prs.save(f)
    print(f"wrote {f}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
