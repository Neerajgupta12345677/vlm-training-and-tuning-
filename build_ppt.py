"""Generate the hackathon submission deck — 5 slides, demo-day standard.

Structure: Title/Problem -> Architecture -> Engineering discipline -> Results
-> Limitations & what's next. One idea per slide, deliberately, rather than
the earlier 2-slide version that packed two ideas onto each.

Built from the MEASURED numbers in PROGRESS.md / HANDOVER.md and the local
arena reproduction, not from claims — every figure here has a command behind
it in the repo.

    python build_ppt.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- palette
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x5B, 0x66, 0x78)
FAINT = RGBColor(0x8A, 0x93, 0xA3)
ACCENT = RGBColor(0xE8, 0x71, 0x22)     # amber — the one accent, spent once per slide
GOOD = RGBColor(0x1B, 0x7F, 0x5A)
BAD = RGBColor(0xB4, 0x33, 0x2E)
RULE = RGBColor(0xD8, 0xDE, 0xE6)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
NAVY = RGBColor(0x0E, 0x1B, 0x2E)       # title slide ground

W, H = Inches(13.333), Inches(7.5)


# ---------------------------------------------------------------- helpers
def _box(slide, x, y, w, h, anchor=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _p(tf, text, size, bold=False, color=INK, space_before=0, space_after=2,
       first=False, align=PP_ALIGN.LEFT, italic=False, font="Segoe UI"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def _kv_line(tf, k, v, size=12, kcolor=INK, vcolor=MUTED, first=False, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(2)
    r = p.add_run(); r.text = k
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = kcolor; r.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = "  " + v
    r2.font.size = Pt(size); r2.font.color.rgb = vcolor; r2.font.name = "Segoe UI"
    return p


def _rect(slide, x, y, w, h, fill, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = RULE
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, kicker, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(0.09), ACCENT)
    tf = _box(slide, Inches(0.6), Inches(0.34), Inches(12.2), Inches(1.05))
    _p(tf, kicker, 11.5, bold=True, color=ACCENT, first=True)
    _p(tf, title, 27, bold=True, color=INK, space_before=1)
    if subtitle:
        _p(tf, subtitle, 12.5, color=MUTED, space_before=3)
    _rect(slide, Inches(0.6), Inches(1.42), Inches(12.13), Inches(0.012), RULE)


def footer(slide, n):
    tf = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf, "FlytBase AHC — Visual Intelligence Hackathon", 9.5, color=FAINT, first=True)
    tf2 = _box(slide, Inches(0.6), Inches(7.16), Inches(12.13), Inches(0.3))
    _p(tf2, str(n), 9.5, color=FAINT, first=True, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ================================================================== SLIDE 1
def slide_title(prs):
    s = new_slide(prs)
    _bg(s, NAVY)
    _rect(s, 0, Inches(6.9), W, Inches(0.6), ACCENT)

    tf = _box(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
    _p(tf, "FLYTBASE AHC · VISUAL INTELLIGENCE HACKATHON", 13, bold=True,
       color=ACCENT, first=True)

    tf2 = _box(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(2.0))
    _p(tf2, "Context-Dependent Anomaly Detection", 40, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF), first=True)
    _p(tf2, "in Aerial Drone Video — Real-Time, on a 4 GB GPU", 40, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF), space_before=2)

    tf3 = _box(s, Inches(0.9), Inches(4.35), Inches(10.6), Inches(1.0))
    _p(tf3, "Noticing is cheap; reasoning is not. So they run as different "
            "stages — a small VLM is consulted on only the 0.19–1.9% of frames "
            "that earn it, and every other decision is made by arithmetic.",
       15.5, color=RGBColor(0xC9, 0xD2, 0xE0), first=True, italic=True)

    # stat strip
    stats = [("50.7 / 100", "public arena · 4th of 7"),
             ("14.9 fps", "full pipeline, real-time"),
             ("GTX 1650", "4 GB, no Tensor Cores"),
             ("3.3 GB", "fine-tuned VLM, merged & live")]
    y = Inches(5.75)
    cw = Inches(2.85)
    for i, (big, small) in enumerate(stats):
        x = Inches(0.9) + i * (cw + Inches(0.12))
        t = _box(s, x, y, cw, Inches(0.9))
        _p(t, big, 22, bold=True, color=ACCENT, first=True)
        _p(t, small, 10.5, color=RGBColor(0x9A, 0xA6, 0xB8), space_before=2)
        if i < 3:
            _rect(s, x + cw + Inches(0.02), y + Inches(0.05), Inches(0.012), Inches(0.7),
                  RGBColor(0x2A, 0x38, 0x50))


# ================================================================== SLIDE 2
def slide_architecture(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "APPROACH", "Three Stages, Split by What Each Component Is Good At",
           "Cheap per-track arithmetic decides the boolean; the VLM only describes what it sees.")

    y = Inches(1.72)
    stages = [
        ("STAGE 1 · every frame", "YOLO26n + ByteTrack",
         "Ego-motion compensated (Shi-Tomasi + LK + RANSAC partial affine) — "
         "a drone pans, so a parked car has non-zero image velocity unless "
         "positions are expressed in a stabilised reference frame."),
        ("STAGE 2 · every frame, no model", "Per-track state machine",
         "Dwell time, speed, zone, neighbour motion. Stopped-vehicle, "
         "wrong-way, congestion, loitering rules — all arithmetic, no "
         "inference cost."),
        ("STAGE 2b · sampled frames", "MobileNetV3-Small",
         "Owns appearance classes only — fire, smoke, flood, spill — because "
         "these ARE visible in a single frame. Threshold 0.72, clip-mean "
         "pooled."),
        ("STAGE 3 · event-triggered only", "Qwen2.5-VL-3B (QLoRA fine-tuned)",
         "Consulted on 0.19–1.9% of frames. Describes the scene; may "
         "ESCALATE a hazard to anomalous, never silently clears a stop the "
         "tracker measured."),
    ]
    bw, gap = Inches(2.92), Inches(0.18)
    for i, (kicker, head, body) in enumerate(stages):
        x = Inches(0.6) + i * (bw + gap)
        _rect(s, x, y, bw, Inches(3.0), CARD, line=True)
        _rect(s, x, y, bw, Inches(0.06), ACCENT if i == 3 else RULE)
        t = _box(s, x + Inches(0.2), y + Inches(0.24), bw - Inches(0.4), Inches(2.6))
        _p(t, kicker, 9.5, bold=True, color=ACCENT if i == 3 else FAINT, first=True)
        _p(t, head, 14.5, bold=True, color=INK, space_before=6)
        _p(t, body, 11, color=MUTED, space_before=8)
        if i < 3:
            a = _box(s, x + bw + Inches(0.01), y + Inches(1.35), gap, Inches(0.3))
            _p(a, "›", 18, bold=True, color=RULE, first=True, align=PP_ALIGN.CENTER)

    # the load-bearing constraint, as a single strip beneath
    y2 = Inches(5.0)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.85), RGBColor(0xFD, 0xF3, 0xE9), line=True)
    _rect(s, Inches(0.6), y2, Inches(0.06), Inches(1.85), ACCENT)
    t = _box(s, Inches(0.95), y2 + Inches(0.2), Inches(11.5), Inches(1.5))
    _p(t, "THE CONSTRAINT EVERYTHING BENDS AROUND", 10.5, bold=True, color=ACCENT, first=True)
    _p(t, "A still frame does not contain motion — a moving car and a stopped car are "
          "pixel-identical in one image. Asked to judge motion directly, Qwen2.5-VL-3B "
          "scored 3/6 (chance) across four prompt revisions, while its prose stayed "
          "accurate. That is why the boolean verdict belongs to Stage 2's arithmetic, "
          "and the VLM's job is narrowed to description plus one-way escalation.",
       12.5, color=INK, space_before=6)
    footer(s, 2)


# ================================================================== SLIDE 3
def slide_discipline(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "ENGINEERING", "What We Did With the Measurements",
           "Three decisions this build made — one shipped, one rejected, one refused.")

    cards = [
        ("SHIPPED", GOOD,
         "Fine-tuned the VLM, merged end-to-end",
         "QLoRA on Kaggle T4 (4 GB VRAM cannot hold weights + gradients + "
         "optimizer state locally) → merged into fp16 → GGUF → Q4_K_M → "
         "Ollama.",
         "Verified 504 LoRA modules injected before export — a silent "
         "no-op merge would have shipped the stock model labelled "
         "\"fine-tuned\". It now answers correctly on a case stock scored "
         "at chance."),
        ("REJECTED", BAD,
         "Killed our own idea with data",
         "Built a second classifier over ego-compensated motion-difference "
         "frames, aimed at the classes the first classifier couldn't "
         "localise.",
         "Trained well (val macro-recall 0.705) — then failed both real "
         "tests: fired on 100% of a clip that is 2% anomalous, and NET −1 "
         "on class accuracy. Not shipped."),
        ("REFUSED", MUTED,
         "Declined to fit the visible test set",
         "Per-class decision thresholds tuned on the public 34-video test "
         "set lift in-sample macro-F1 to 0.289.",
         "Its own leave-one-video-out check falls to 0.230 — worse than "
         "the untuned global rule. Three separate threshold-fitting "
         "attempts measured and rejected on this basis."),
    ]
    y = Inches(1.72)
    cw = Inches(3.86)
    for i, (tag, tagcol, head, body1, body2) in enumerate(cards):
        x = Inches(0.6) + i * (cw + Inches(0.2))
        _rect(s, x, y, cw, Inches(4.9), CARD, line=True)
        _rect(s, x, y, cw, Inches(0.42), tagcol)
        tg = _box(s, x, y, cw, Inches(0.42))
        _p(tg, tag, 12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), first=True,
           align=PP_ALIGN.CENTER)
        t = _box(s, x + Inches(0.24), y + Inches(0.62), cw - Inches(0.48), Inches(4.1))
        _p(t, head, 15.5, bold=True, color=INK, first=True)
        _p(t, body1, 11.5, color=INK, space_before=10)
        _p(t, body2, 11.5, color=MUTED, space_before=10)
    footer(s, 3)


# ================================================================== SLIDE 4
def slide_results(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "RESULTS", "Measured, Not Claimed",
           "Every number below has a reproduction command in the repo.")

    y = Inches(1.72)
    # throughput card
    _rect(s, Inches(0.6), y, Inches(5.95), Inches(2.7), CARD, line=True)
    t = _box(s, Inches(0.9), y + Inches(0.22), Inches(5.4), Inches(2.3))
    _p(t, "THROUGHPUT · 4K 25 FPS SOURCE, GTX 1650", 10.5, bold=True, color=MUTED, first=True)
    perf = [("Stage 1, threaded decode", "26.6 fps", "1.06× real-time on 4K", GOOD),
            ("Full pipeline, rules, stride 2", "14.9 fps", "vs 12.5 needed — REAL-TIME", GOOD),
            ("Serial decode (reverted)", "20.0 fps", "decode + detect serialise", MUTED),
            ("Hybrid, VLM in the loop", "6.8 fps", "one 4GB card can't run both", BAD)]
    for i, (a, b, c, col) in enumerate(perf):
        yy = y + Inches(0.56) + Inches(0.46) * i
        tr = _box(s, Inches(0.9), yy, Inches(5.4), Inches(0.44))
        p = tr.paragraphs[0]
        r = p.add_run(); r.text = a
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = INK; r.font.name = "Segoe UI"
        r = p.add_run(); r.text = "   " + b
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = col; r.font.name = "Consolas"
        p2 = tr.add_paragraph()
        r = p2.add_run(); r.text = c
        r.font.size = Pt(10); r.font.color.rgb = MUTED; r.font.name = "Segoe UI"

    # scoreboard card
    _rect(s, Inches(6.78), y, Inches(5.95), Inches(2.7), CARD, line=True)
    t2 = _box(s, Inches(7.08), y + Inches(0.22), Inches(5.4), Inches(2.3))
    _p(t2, "PUBLIC ARENA · 50.7 / 100 · 4TH OF 7", 10.5, bold=True, color=ACCENT, first=True)
    board = [("Level 1  clip classification", "17.5 / 25", GOOD),
             ("Level 2  when it happens", "25.2 / 35", GOOD),
             ("Level 3  long context", "8.0 / 40", BAD)]
    for i, (a, b, col) in enumerate(board):
        yy = y + Inches(0.6) + Inches(0.44) * i
        tr = _box(s, Inches(7.08), yy, Inches(5.4), Inches(0.4))
        p = tr.paragraphs[0]
        r = p.add_run(); r.text = a
        r.font.size = Pt(12.5); r.font.color.rgb = INK; r.font.name = "Segoe UI"
        r = p.add_run(); r.text = "     " + b
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = col; r.font.name = "Consolas"
    _p(_box(s, Inches(7.08), y + Inches(2.0), Inches(5.4), Inches(0.6)),
       "score_arena.py reproduces the live leaderboard exactly — every match "
       "and false-alarm count — so changes are measured locally, never guessed.",
       10, color=MUTED, first=True)

    # detector table, full width
    y3 = Inches(4.62)
    _rect(s, Inches(0.6), y3, Inches(12.13), Inches(2.15), CARD, line=True)
    t3 = _box(s, Inches(0.9), y3 + Inches(0.2), Inches(11.5), Inches(1.8))
    _p(t3, "THE DETECTOR WAS THE ORIGINAL WEAK LINK · VisDrone val, 100 images, 5,189 objects",
       10.5, bold=True, color=MUTED, first=True)
    rows = [("config", "overall recall", "tiny objects", True),
            ("stock — imgsz 640, conf 0.25", "0.152", "0.008", False),
            ("aerial — imgsz 1280, conf 0.10", "0.413", "0.149", False)]
    ty = y3 + Inches(0.66)
    for i, (c1, c2, c3, hdr) in enumerate(rows):
        yy = ty + Inches(0.34) * i
        tr = _box(s, Inches(0.9), yy, Inches(8.5), Inches(0.32))
        p = tr.paragraphs[0]
        for txt, wdt, bold, col in ((c1, 38, hdr or i == 2, MUTED if hdr else INK),
                                    (c2, 18, hdr or i == 2, MUTED if hdr else (GOOD if i == 2 else INK)),
                                    (c3, 16, hdr or i == 2, MUTED if hdr else (GOOD if i == 2 else INK))):
            r = p.add_run(); r.text = txt.ljust(wdt)
            r.font.size = Pt(12); r.font.bold = bold; r.font.color.rgb = col
            r.font.name = "Consolas"
    _p(_box(s, Inches(9.6), ty + Inches(0.02), Inches(3.0), Inches(1.0)),
       "Stock settings miss 85% of aerial objects. Config alone nearly triples recall, "
       "for free, and still clears 25 fps.", 10.5, color=MUTED, first=True)
    footer(s, 4)


# ================================================================== SLIDE 5
def slide_forward(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "LIMITATIONS & WHAT'S NEXT", "The Open Problem Is Named, Not Hidden",
           "Level 3 is a diagnosed gap with a specific next step — not an unknown.")

    y = Inches(1.72)
    _rect(s, Inches(0.6), y, Inches(7.0), Inches(4.75), CARD, line=True)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(4.75), BAD)
    t = _box(s, Inches(0.95), y + Inches(0.24), Inches(6.4), Inches(4.3))
    _p(t, "WHY LEVEL 3 SCORES 8.0 / 40", 11, bold=True, color=BAD, first=True)
    _p(t, "The evidence is a property of one track, not of any frame.", 16, bold=True,
       color=INK, space_before=6)
    _p(t, "On a 6-minute clip, 2% of the runtime is anomalous (a person "
          "dwelling too long) while 98% shows ordinary motion in the same "
          "scene. Any whole-frame classifier — RGB or motion-difference — "
          "answers \"what kind of scene is this\", which cannot separate "
          "that 2% from the rest, at any threshold.", 12.5, color=MUTED, space_before=8)
    _p(t, "Measured, not assumed:", 12, bold=True, color=INK, space_before=12)
    facts = [
        "Confidence-ranked windows put the wrong one first in all 4 L3 videos.",
        "A dedicated motion-difference classifier still fired on 100% of a "
        "clip whose ground truth is 2% anomalous.",
        "Stage 2's tracker already computes per-track dwell, speed and zone "
        "— the only component with genuine temporal state.",
    ]
    for f in facts:
        p = t.add_paragraph(); p.space_before = Pt(5)
        r = p.add_run(); r.text = "▪ "
        r.font.size = Pt(11.5); r.font.color.rgb = BAD; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = f
        r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"

    # right column: next steps + honesty note
    x2 = Inches(7.86)
    _rect(s, x2, y, Inches(4.87), Inches(2.85), CARD, line=True)
    t2 = _box(s, x2 + Inches(0.28), y + Inches(0.22), Inches(4.35), Inches(2.5))
    _p(t2, "NEXT STEP", 11, bold=True, color=ACCENT, first=True)
    _p(t2, "Emit intervals from the tracker's own state, not from a "
          "classifier's per-frame score.", 13, bold=True, color=INK, space_before=6)
    for item in ["Stage 2 already flags dwell-time crossings per track",
                 "Convert those crossings directly into timed windows",
                 "Reserve the classifier for what it's proven at: "
                 "Level 1/2 class labels"]:
        p = t2.add_paragraph(); p.space_before = Pt(6)
        r = p.add_run(); r.text = "→ "
        r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"

    y2 = y + Inches(3.05)
    _rect(s, x2, y2, Inches(4.87), Inches(1.7), RGBColor(0xF0, 0xF3, 0xF7), line=True)
    t3 = _box(s, x2 + Inches(0.28), y2 + Inches(0.2), Inches(4.35), Inches(1.35))
    _p(t3, "WHY WE'RE SHOWING THIS", 10, bold=True, color=MUTED, first=True)
    _p(t3, "A negative result, measured and understood, is worth more than a "
          "vague claim of success. Every rejection on this deck cost a real "
          "GPU run to disprove — none was assumed.", 11.5, color=INK, space_before=6)
    footer(s, 5)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_title(prs)
    slide_architecture(prs)
    slide_discipline(prs)
    slide_results(prs)
    slide_forward(prs)
    out = Path("deck")
    out.mkdir(exist_ok=True)
    f = out / "FlytBase_AHC_Submission.pptx"
    try:
        prs.save(f)
    except PermissionError:
        # The previous file is likely open in PowerPoint and locked - write
        # alongside it rather than fail the whole build.
        f = out / "FlytBase_AHC_Submission_v2.pptx"
        prs.save(f)
    print(f"wrote {f}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
