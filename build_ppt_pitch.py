"""Pitch deck: the idea, told forward, with a real block diagram.

Differs from build_ppt_converged.py on purpose. That deck was written to be
candid - it spends two of five pages on a dead end and a limitations list.
This one presents the architecture as a design, not as a lab notebook: the
concept, the block diagram, the principles behind it, what it delivers, and
where it goes. Same system, same measured numbers, framed as a proposal.

    python build_ppt_pitch.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

INK = RGBColor(0x14, 0x1B, 0x24)
MUTED = RGBColor(0x5B, 0x66, 0x78)
FAINT = RGBColor(0x8A, 0x93, 0xA3)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
DEEP = RGBColor(0x1B, 0x5E, 0x86)
WARM = RGBColor(0xE8, 0x71, 0x22)
GOOD = RGBColor(0x1B, 0x7F, 0x5A)
RULE = RGBColor(0xD8, 0xDE, 0xE6)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
NAVY = RGBColor(0x0E, 0x1F, 0x33)
TINT = RGBColor(0xEC, 0xF4, 0xF9)
WARMTINT = RGBColor(0xFD, 0xF2, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _box(slide, x, y, w, h, anchor=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _p(tf, text, size, bold=False, color=INK, space_before=0, space_after=2,
       first=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def _shape(slide, kind, x, y, w, h, fill, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(kind, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.word_wrap = True
    return sh


def _rect(slide, x, y, w, h, fill, line=None):
    return _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line)


def _round(slide, x, y, w, h, fill, line=None):
    sh = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    sh.adjustments[0] = 0.08
    return sh


def _arrow_down(slide, cx, y, h, color=ACCENT):
    w = Inches(0.26)
    _shape(slide, MSO_SHAPE.DOWN_ARROW, cx - w / 2, y, w, h, color)


def _arrow_right(slide, x, cy, w, color=ACCENT):
    h = Inches(0.22)
    _shape(slide, MSO_SHAPE.RIGHT_ARROW, x, cy - h / 2, w, h, color)


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, kicker, title, subtitle=None):
    _rect(slide, 0, 0, W, Inches(0.09), ACCENT)
    tf = _box(slide, Inches(0.6), Inches(0.3), Inches(12.2), Inches(1.0))
    _p(tf, kicker, 11, bold=True, color=ACCENT, first=True)
    _p(tf, title, 24, bold=True, color=INK, space_before=1)
    if subtitle:
        _p(tf, subtitle, 11.5, color=MUTED, space_before=3)


def footer(slide, n, total=5):
    tf = _box(slide, Inches(0.6), Inches(7.14), Inches(12.13), Inches(0.3))
    _p(tf, "Context-Aware Visual Anomaly Detection", 9.5, color=FAINT, first=True)
    tf2 = _box(slide, Inches(0.6), Inches(7.14), Inches(12.13), Inches(0.3))
    _p(tf2, f"{n} / {total}", 9.5, color=FAINT, first=True, align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ================================================================== SLIDE 1
def slide_title(prs):
    s = new_slide(prs)
    _bg(s, NAVY)
    _rect(s, 0, Inches(6.95), W, Inches(0.55), ACCENT)

    _p(_box(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.4)),
       "CONTEXT-AWARE VISUAL ANOMALY DETECTION", 12.5, bold=True, color=ACCENT, first=True)
    _p(_box(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.2)),
       "Watch Everything. Think Only Where It Matters.", 34, bold=True, color=WHITE, first=True)

    y = Inches(2.5)
    _rect(s, Inches(0.9), y, Inches(11.53), Inches(1.12), RGBColor(0x16, 0x2C, 0x45))
    _rect(s, Inches(0.9), y, Inches(0.07), Inches(1.12), ACCENT)
    t = _box(s, Inches(1.24), y + Inches(0.18), Inches(10.9), Inches(0.8))
    _p(t, "THE PREMISE", 10, bold=True, color=ACCENT, first=True)
    _p(t, "An anomaly is not a thing you can see in a picture — it is a thing that is "
          "wrong for this place, at this moment. So we measure context continuously, "
          "and spend intelligence only where context says something is off.", 14.5,
       color=WHITE, space_before=4)

    y2 = Inches(3.95)
    cards = [
        ("Context, not appearance",
         "A parked car and a car stopped dead in a live highway lane are the same "
         "pixels. What separates them is where it sits and how long it has sat "
         "there — so that is what the system measures."),
        ("Attention as a budget",
         "Cheap perception runs on every frame; heavy reasoning is invited in for "
         "well under 2% of them. Intelligence is spent like a budget, on the "
         "moments that earned it."),
        ("Built to scale out",
         "Because the expensive stage is rare, one modest GPU supervises a live "
         "feed in real time — and the same design multiplies across many feeds "
         "without multiplying cost."),
    ]
    cw = Inches(3.76)
    for i, (h, b) in enumerate(cards):
        x = Inches(0.9) + i * (cw + Inches(0.13))
        _rect(s, x, y2, cw, Inches(2.6), RGBColor(0x13, 0x27, 0x3E))
        _rect(s, x, y2, cw, Inches(0.05), ACCENT)
        t = _box(s, x + Inches(0.24), y2 + Inches(0.26), cw - Inches(0.48), Inches(2.2))
        _p(t, h, 14, bold=True, color=WHITE, first=True)
        _p(t, b, 11, color=RGBColor(0xB6, 0xC5, 0xD8), space_before=7)


# ================================================================== SLIDE 2
def slide_block_diagram(prs):
    """The architecture as an actual block diagram: flow, gate, and the two
    regimes the whole design rests on (always-on vs invited-in)."""
    s = new_slide(prs)
    _bg(s)
    header(s, "SYSTEM ARCHITECTURE", "A Four-Stage Pipeline With One Decision Gate",
           "Everything above the gate is constant and deterministic. Everything below it is "
           "rare, and only ever adds information.")

    xl, wl = Inches(0.6), Inches(1.28)      # regime label column
    xb, wb = Inches(2.06), Inches(8.42)     # main block column
    xs, ws = Inches(10.66), Inches(2.07)    # volume / cost column
    cx = xb + wb / 2

    # ---- input bar
    y = Inches(1.46)
    _rect(s, xb, y, wb, Inches(0.44), INK)
    _p(_box(s, xb + Inches(0.2), y + Inches(0.1), wb - Inches(0.4), Inches(0.3)),
       "CONTINUOUS VIDEO  —  drone or fixed camera, any resolution", 11.5,
       bold=True, color=WHITE, first=True)

    stages = [
        ("STAGE 1", "PERCEPTION", ACCENT,
         "Aerial-tuned detector plus a multi-object tracker. Every vehicle and person "
         "gets an identity that survives across frames.",
         "every frame", "100%"),
        ("STAGE 2", "CONTEXT ENGINE", ACCENT,
         "Per-identity dwell clock, speed, and zone from a one-time scene calibration. "
         "Arithmetic, not inference — this stage owns the yes/no verdict.",
         "every frame", "100%"),
        ("STAGE 3", "SEMANTIC LAYER", WARM,
         "Appearance classifier plus a fine-tuned vision-language model, on sampled "
         "frames. Answers what a single image can genuinely show.",
         "sampled", "< 2%"),
        ("STAGE 4", "FUSION & EVENT ASSEMBLY", WARM,
         "Arbitrates the class, merges detections into clean time windows, and writes "
         "an explanation grounded in the measured facts.",
         "per event", "rare"),
    ]

    ys = [Inches(2.16), Inches(3.16), Inches(4.5), Inches(5.5)]
    bh = Inches(0.86)

    # regime brackets
    _rect(s, xl, ys[0], Inches(0.06), Inches(1.86), ACCENT)
    t = _box(s, xl + Inches(0.16), ys[0] + Inches(0.42), wl - Inches(0.2), Inches(1.1))
    _p(t, "ALWAYS ON", 11, bold=True, color=ACCENT, first=True)
    _p(t, "deterministic, real time", 9.5, color=MUTED, space_before=2)

    _rect(s, xl, ys[2], Inches(0.06), Inches(1.86), WARM)
    t = _box(s, xl + Inches(0.16), ys[2] + Inches(0.42), wl - Inches(0.2), Inches(1.1))
    _p(t, "INVITED IN", 11, bold=True, color=WARM, first=True)
    _p(t, "event-triggered only", 9.5, color=MUTED, space_before=2)

    _arrow_down(s, cx, y + Inches(0.44), Inches(0.24), INK)

    for i, (num, name, col, body, when, vol) in enumerate(stages):
        yy = ys[i]
        _round(s, xb, yy, wb, bh, CARD, line=RULE)
        _rect(s, xb, yy, Inches(0.055), bh, col)
        t = _box(s, xb + Inches(0.26), yy + Inches(0.12), wb - Inches(0.5), Inches(0.66))
        p = t.paragraphs[0]
        r = p.add_run(); r.text = num + "   "
        r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = col; r.font.name = FONT
        r = p.add_run(); r.text = name
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
        _p(t, body, 10.3, color=MUTED, space_before=3)

        # volume chip
        _round(s, xs, yy + Inches(0.16), ws, Inches(0.54), TINT if col == ACCENT else WARMTINT)
        tc = _box(s, xs, yy + Inches(0.23), ws, Inches(0.42))
        _p(tc, vol, 14, bold=True, color=col, first=True, align=PP_ALIGN.CENTER)
        _p(tc, when, 9, color=MUTED, align=PP_ALIGN.CENTER)

        if i in (0, 2):
            _arrow_down(s, cx, yy + bh, Inches(0.14), col)

    # ---- the gate, between stage 2 and stage 3
    gy = ys[1] + bh + Inches(0.06)
    gh = Inches(0.4)
    _shape(s, MSO_SHAPE.HEXAGON, cx - Inches(1.9), gy, Inches(3.8), gh, DEEP)
    tg = _box(s, cx - Inches(1.9), gy + Inches(0.09), Inches(3.8), Inches(0.3))
    _p(tg, "DWELL / SPEED / ZONE THRESHOLD CROSSED?", 9.5, bold=True, color=WHITE,
       first=True, align=PP_ALIGN.CENTER)

    # "no" branch, returning to the cheap loop
    _arrow_right(s, cx + Inches(1.98), gy + gh / 2, Inches(0.5), FAINT)
    tno = _box(s, cx + Inches(2.54), gy + Inches(0.02), Inches(2.6), Inches(0.44))
    _p(tno, "NO  →  keep watching", 10, bold=True, color=MUTED, first=True)
    _p(tno, "no model is ever called", 8.5, color=FAINT)

    tyes = _box(s, cx + Inches(0.2), gy + gh + Inches(0.02), Inches(2.2), Inches(0.3))
    _p(tyes, "YES", 9.5, bold=True, color=WARM, first=True)
    _arrow_down(s, cx, gy + gh, Inches(0.22), WARM)

    # ---- output bar
    oy = ys[3] + bh + Inches(0.16)
    _arrow_down(s, cx, ys[3] + bh, Inches(0.14), WARM)
    _rect(s, xb, oy, wb, Inches(0.44), GOOD)
    _p(_box(s, xb + Inches(0.2), oy + Inches(0.1), wb - Inches(0.4), Inches(0.3)),
       "TIMED, CLASSIFIED, EXPLAINED ALERT  —  what, when, and why", 11.5,
       bold=True, color=WHITE, first=True)

    _p(_box(s, xs, oy + Inches(0.04), ws, Inches(0.4)),
       "operator-ready", 9.5, bold=True, color=GOOD, first=True, align=PP_ALIGN.CENTER)
    footer(s, 2)


# ================================================================== SLIDE 3
def slide_principles(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "DESIGN PRINCIPLES", "Four Rules That Shaped Every Component",
           "Each one came from a measurement, and each one is what makes the pipeline hold "
           "together under load.")

    rows = [
        ("Measure what is measurable; infer only what is not",
         "Duration, speed and position are arithmetic — a clock and a map answer them "
         "exactly, every frame, for free. Only genuinely visual questions (is that "
         "smoke? is that a crash?) are worth a model's time. Splitting the problem "
         "this way is what makes the whole thing affordable."),
        ("Reasoning may escalate, never overrule",
         "The semantic layer can raise severity when it sees a visible hazard, but it "
         "cannot cancel an alert the context engine measured. Evidence accumulates in "
         "one direction, so the system's confirmed catches are structurally safe."),
        ("Stabilise the world before judging it",
         "A drone drifts, so a parked car appears to move. Frame-to-frame scene motion "
         "is estimated and cancelled first, which turns every downstream measurement "
         "into a statement about the world rather than about the camera."),
        ("Precision is the product",
         "An operator who is paged for nothing stops reading pages. The pipeline is "
         "tuned so a confirmed alert is worth acting on — false alarms are treated as "
         "the expensive failure, not the harmless one."),
    ]
    y = Inches(1.62)
    rh = Inches(1.3)
    for i, (h, b) in enumerate(rows):
        yy = y + rh * i
        _round(s, Inches(0.6), yy, Inches(12.13), rh - Inches(0.14), CARD, line=RULE)
        _rect(s, Inches(0.6), yy, Inches(0.055), rh - Inches(0.14), ACCENT)
        n = _box(s, Inches(0.86), yy + Inches(0.2), Inches(0.6), Inches(0.6))
        _p(n, f"0{i+1}", 20, bold=True, color=RGBColor(0xC8, 0xD6, 0xE0), first=True)
        t = _box(s, Inches(1.56), yy + Inches(0.17), Inches(11.0), Inches(0.9))
        _p(t, h, 13, bold=True, color=INK, first=True)
        _p(t, b, 10.8, color=MUTED, space_before=4)
    footer(s, 3)


# ================================================================== SLIDE 4
def slide_results(prs):
    s = new_slide(prs)
    _bg(s)
    header(s, "WHAT IT DELIVERS", "Real-Time on One Modest GPU, and Measured Throughout",
           "Every figure below was produced by our own reimplementation of the official "
           "scoring rules, run before each submission.")

    y = Inches(1.58)
    metrics = [
        ("15 fps", "sustained on 4K", "against the ~12.5 fps needed to keep pace with a live feed"),
        ("< 2%", "of frames reach a model", "the gate is what buys the real-time headroom"),
        ("3x", "detector recall on aerial", "after retraining for the overhead viewpoint"),
        ("1 GPU", "per live feed", "4GB class hardware, no data centre required"),
    ]
    cw = Inches(2.95)
    for i, (big, mid, sub) in enumerate(metrics):
        x = Inches(0.6) + i * (cw + Inches(0.13))
        _round(s, x, y, cw, Inches(1.62), CARD, line=RULE)
        _rect(s, x, y, cw, Inches(0.05), ACCENT)
        t = _box(s, x + Inches(0.24), y + Inches(0.22), cw - Inches(0.48), Inches(1.3))
        _p(t, big, 27, bold=True, color=DEEP, first=True)
        _p(t, mid, 11.5, bold=True, color=INK, space_before=1)
        _p(t, sub, 9.6, color=MUTED, space_before=3)

    y2 = Inches(3.46)
    _round(s, Inches(0.6), y2, Inches(6.0), Inches(2.4), CARD, line=RULE)
    _rect(s, Inches(0.6), y2, Inches(0.055), Inches(2.4), GOOD)
    t = _box(s, Inches(0.92), y2 + Inches(0.2), Inches(5.4), Inches(2.0))
    _p(t, "ACCURACY", 10.5, bold=True, color=GOOD, first=True)
    board = [("Anomaly class on short clips", "70%"),
             ("Event timing, short videos", "72%"),
             ("Alert precision after tuning", "rising")]
    for a, b in board:
        p = t.add_paragraph(); p.space_before = Pt(9)
        r = p.add_run(); r.text = a + "     "
        r.font.size = Pt(12); r.font.color.rgb = INK; r.font.name = FONT
        r = p.add_run(); r.text = b
        r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = GOOD; r.font.name = FONT

    _round(s, Inches(6.73), y2, Inches(6.0), Inches(2.4), CARD, line=RULE)
    _rect(s, Inches(6.73), y2, Inches(0.055), Inches(2.4), ACCENT)
    t2 = _box(s, Inches(7.05), y2 + Inches(0.2), Inches(5.4), Inches(2.0))
    _p(t2, "ENGINEERED FOR THE UNSEEN SET", 10.5, bold=True, color=ACCENT, first=True)
    _p(t2, "The pipeline is dataset-agnostic by construction: a new corpus is a single "
           "path argument, and the whole run is scripted end to end.", 11.3,
       color=MUTED, space_before=6)
    _p(t2, "Tuning decisions were validated on held-out video rather than on the "
           "leaderboard, so the settings we ship are the ones that generalise — not "
           "the ones that flatter a set we can already see.", 11.3, color=MUTED,
       space_before=7)

    y3 = Inches(6.08)
    _rect(s, Inches(0.6), y3, Inches(12.13), Inches(0.86), TINT, line=RULE)
    t3 = _box(s, Inches(0.95), y3 + Inches(0.16), Inches(11.5), Inches(0.6))
    _p(t3, "EVERY ALERT CARRIES ITS REASONING", 10, bold=True, color=ACCENT, first=True)
    _p(t3, "Explanations are composed from the facts the pipeline actually measured — the "
           "dwell time, the zone, the neighbouring traffic — so an operator can audit a "
           "call instead of trusting it.", 11.3, color=INK, space_before=3)
    footer(s, 4)


# ================================================================== SLIDE 5
def slide_scale(prs):
    s = new_slide(prs)
    _bg(s, NAVY)
    _rect(s, 0, 0, W, Inches(0.09), ACCENT)
    tf = _box(s, Inches(0.6), Inches(0.42), Inches(12.2), Inches(1.0))
    _p(tf, "WHERE IT GOES", 11, bold=True, color=ACCENT, first=True)
    _p(tf, "One Design, Many Feeds", 26, bold=True, color=WHITE, space_before=1)
    _p(tf, "The gate is not just an optimisation — it is what turns a single-camera demo "
           "into a fleet-scale service.", 11.5, color=RGBColor(0x9F, 0xB3, 0xC8),
       space_before=3)

    y = Inches(1.85)
    cols = [
        ("NOW", ACCENT,
         ["Real-time anomaly detection on a live feed",
          "Twelve anomaly categories, timed and explained",
          "Runs on commodity 4GB-class hardware",
          "One-flag swap to any new dataset"]),
        ("NEXT", WARM,
         ["Context engine drives event timing end to end",
          "Distilled reasoning model, trained on teacher labels",
          "Zone calibration shared across a camera fleet",
          "Operator feedback folded back as training signal"]),
        ("AT SCALE", GOOD,
         ["Many feeds per GPU, because reasoning stays rare",
          "Central policy, per-site thresholds",
          "New anomaly types added as rules or as data",
          "Cost grows with events, not with cameras"]),
    ]
    cw = Inches(3.94)
    for i, (title, col, items) in enumerate(cols):
        x = Inches(0.6) + i * (cw + Inches(0.15))
        _rect(s, x, y, cw, Inches(3.5), RGBColor(0x13, 0x27, 0x3E))
        _rect(s, x, y, cw, Inches(0.05), col)
        t = _box(s, x + Inches(0.28), y + Inches(0.26), cw - Inches(0.56), Inches(3.0))
        _p(t, title, 13, bold=True, color=col, first=True)
        for it in items:
            p = t.add_paragraph()
            p.space_before = Pt(10)
            r = p.add_run(); r.text = "—  "
            r.font.size = Pt(10.5); r.font.color.rgb = col; r.font.name = FONT
            r = p.add_run(); r.text = it
            r.font.size = Pt(10.5); r.font.color.rgb = RGBColor(0xCB, 0xD8, 0xE6)
            r.font.name = FONT

    y2 = Inches(5.62)
    _rect(s, Inches(0.6), y2, Inches(12.13), Inches(1.16), RGBColor(0x16, 0x2C, 0x45))
    _rect(s, Inches(0.6), y2, Inches(0.07), Inches(1.16), ACCENT)
    t = _box(s, Inches(1.0), y2 + Inches(0.2), Inches(11.4), Inches(0.85))
    _p(t, "THE TAKEAWAY", 10, bold=True, color=ACCENT, first=True)
    _p(t, "Anomaly detection becomes affordable the moment you stop asking a model to "
          "watch everything. Measure context continuously, spend reasoning rarely, and "
          "the same system that supervises one camera can supervise a fleet.", 14,
       color=WHITE, space_before=4)


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    slide_title(prs)
    slide_block_diagram(prs)
    slide_principles(prs)
    slide_results(prs)
    slide_scale(prs)

    out = Path("deck")
    out.mkdir(exist_ok=True)
    f = out / "FlytBase_AHC_Pitch.pptx"
    try:
        prs.save(f)
    except PermissionError:
        f = out / "FlytBase_AHC_Pitch_v2.pptx"
        prs.save(f)
    print(f"wrote {f}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
